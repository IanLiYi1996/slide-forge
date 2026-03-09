"""Package slide images into a .pptx presentation.

Usage:
    python generate_pptx.py <image_dir> <output_pptx> [--title "Presentation Title"]

The script reads all PNG/JPG images from <image_dir> in sorted order
and creates a 16:9 widescreen presentation with each image as a full-page slide.
"""

import argparse
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Emu


def create_pptx(image_dir: Path, output_path: Path, title: str = "Presentation") -> None:
    """Create a .pptx file from a directory of slide images."""
    # Collect image files sorted by name
    image_extensions = {".png", ".jpg", ".jpeg", ".webp"}
    images = sorted(
        [f for f in image_dir.iterdir() if f.suffix.lower() in image_extensions],
        key=lambda f: f.name,
    )

    if not images:
        print(f"Error: No images found in {image_dir}", file=sys.stderr)
        sys.exit(1)

    print(f"Found {len(images)} slide images")

    # Create presentation with 16:9 widescreen dimensions
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Use blank layout
    blank_layout = prs.slide_layouts[6]  # Blank layout

    for i, img_path in enumerate(images):
        slide = prs.slides.add_slide(blank_layout)

        # Add image as full-page background
        slide.shapes.add_picture(
            str(img_path),
            left=Emu(0),
            top=Emu(0),
            width=prs.slide_width,
            height=prs.slide_height,
        )

        print(f"  Added slide {i + 1}: {img_path.name}")

    # Save
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))
    print(f"\nSaved: {output_path}")
    print(f"Total slides: {len(images)}")


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Package slide images into .pptx")
    parser.add_argument("image_dir", type=Path, help="Directory containing slide images")
    parser.add_argument("output_pptx", type=Path, help="Output .pptx file path")
    parser.add_argument("--title", default="Presentation", help="Presentation title")
    args = parser.parse_args()

    if not args.image_dir.is_dir():
        print(f"Error: {args.image_dir} is not a directory", file=sys.stderr)
        sys.exit(1)

    create_pptx(args.image_dir, args.output_pptx, args.title)
